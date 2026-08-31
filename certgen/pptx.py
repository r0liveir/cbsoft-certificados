from __future__ import annotations

from pathlib import Path

from .core import CertificateError, TAG_PATTERN, extract_variables


def _text_frames(shape):
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from _text_frames(child)


def _slide_text_frames(slide):
    for shape in slide.shapes:
        yield from _text_frames(shape)


def template_variables(template: Path, slide_number: int) -> set[str]:
    presentation = _load(template)
    slide = _slide(presentation, slide_number)
    return extract_variables(frame.text for frame in _slide_text_frames(slide))


def _load(template: Path):
    try:
        from pptx import Presentation
    except ImportError as error:
        raise CertificateError("Install dependencies first: python -m pip install -e .") from error
    try:
        return Presentation(template)
    except Exception as error:
        raise CertificateError(f"Could not open PowerPoint template {template}.") from error


def _slide(presentation, slide_number: int):
    if not 1 <= slide_number <= len(presentation.slides):
        raise CertificateError(f"Slide {slide_number} does not exist (template has {len(presentation.slides)} slides).")
    return presentation.slides[slide_number - 1]


def _keep_only_slide(presentation, selected_index: int) -> None:
    for index in reversed(range(len(presentation.slides))):
        if index == selected_index:
            continue
        slide_id = presentation.slides._sldIdLst[index]  # python-pptx public API has no deletion helper.
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[index]


def _replace_text(frame, context: dict[str, str]) -> None:
    for paragraph in frame.paragraphs:
        original = "".join(run.text for run in paragraph.runs)
        replaced = TAG_PATTERN.sub(lambda match: context.get(match.group(1), match.group(0)), original)
        if original == replaced:
            continue
        # Normal templates keep each tag in one run, preserving the designer's typography.
        for run in paragraph.runs:
            run.text = TAG_PATTERN.sub(lambda match: context.get(match.group(1), match.group(0)), run.text)
        if TAG_PATTERN.search("".join(run.text for run in paragraph.runs)):
            # A tag split across PowerPoint runs cannot preserve each run's formatting.
            # This fallback keeps the paragraph formatting and makes the certificate usable.
            paragraph.text = replaced


def render(template: Path, slide_number: int, context: dict[str, str], destination: Path) -> None:
    presentation = _load(template)
    selected = _slide(presentation, slide_number)
    for frame in _slide_text_frames(selected):
        _replace_text(frame, context)
    _keep_only_slide(presentation, slide_number - 1)
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(destination)
